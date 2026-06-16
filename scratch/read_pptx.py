from pptx import Presentation

prs = Presentation(r'd:\MAYCON\AGENTES\support\Banner Workshop IA.pptx')

print("Num slides:", len(prs.slides))
print("Width:", round(prs.slide_width.cm, 2), "cm  Height:", round(prs.slide_height.cm, 2), "cm")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    for shape in slide.shapes:
        stype = shape.shape_type
        sname = shape.name
        print(f"  [Shape] name={sname!r}  type={stype}")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    # also capture font size / bold from first run
                    extra = ""
                    if para.runs:
                        r = para.runs[0]
                        sz  = r.font.size.pt if r.font.size else None
                        bold = r.font.bold
                        extra = f"  [size={sz}, bold={bold}]"
                    print(f"    TEXT: {line!r}{extra}")
    print()
